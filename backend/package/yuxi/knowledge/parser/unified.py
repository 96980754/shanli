"""Unified document parsers that produce Markdown and source metadata."""

from __future__ import annotations

import asyncio
import base64
import os
import re
import tempfile
import threading
import time
from dataclasses import dataclass, field
from importlib.metadata import PackageNotFoundError, version
from pathlib import Path
from typing import Any

import aiofiles
import fitz
from docling.datamodel.base_models import InputFormat
from docling.document_converter import DocumentConverter
from langchain_community.document_loaders import PyPDFLoader
from markdownify import markdownify as md_convert

from yuxi.knowledge.parser.ocr_routing import (
    OCRRouteResult,
    OCRRoutingError,
    OCRRoutingPolicy,
    assess_text_quality,
    looks_like_structured_layout,
    notify_processing_stage,
    run_ocr_fallback,
)
from yuxi.knowledge.parser.image_normalization import (
    ImageNormalizationError,
    get_image_format_capability,
    normalize_image_for_ocr,
    validate_image_bytes,
)
from yuxi.knowledge.parser.legacy_office import (
    LEGACY_OFFICE_FORMATS,
    LegacyOfficeConversionError,
    LegacyOfficeConverter,
    get_legacy_office_capability,
    validate_legacy_office_bytes,
    validate_ooxml_bytes,
)
from yuxi.knowledge.parser.zip_utils import process_zip_file as _process_zip_file
from yuxi.storage.minio import get_minio_client
from yuxi.utils import logger

SUPPORTED_FILE_EXTENSIONS: tuple[str, ...] = (
    ".txt",
    ".md",
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    ".doc",
    ".xls",
    ".ppt",
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".webp",
)


def is_supported_file_extension(file_name: str | os.PathLike[str]) -> bool:
    """Return whether the runtime can currently ingest the file extension."""
    suffix = Path(file_name).suffix.lower()
    capability = get_file_format_capability(suffix)
    return bool(capability and capability["enabled"])


def get_file_format_capability(
    suffix: str,
    params: dict[str, Any] | None = None,
) -> dict[str, Any] | None:
    normalized_suffix = suffix.lower()
    if normalized_suffix not in SUPPORTED_FILE_EXTENSIONS:
        return None
    if normalized_suffix in LEGACY_OFFICE_FORMATS:
        return get_legacy_office_capability(normalized_suffix, params=params)
    if normalized_suffix in {".gif", ".webp"}:
        return get_image_format_capability(normalized_suffix)
    return {
        "extension": normalized_suffix,
        "enabled": True,
        "requires_converter": False,
        "availability": "available",
        "reason": None,
    }


def get_file_format_capabilities(params: dict[str, Any] | None = None) -> list[dict[str, Any]]:
    return [
        capability
        for suffix in SUPPORTED_FILE_EXTENSIONS
        if (capability := get_file_format_capability(suffix, params=params)) is not None
    ]


def get_enabled_file_extensions(params: dict[str, Any] | None = None) -> tuple[str, ...]:
    return tuple(
        capability["extension"] for capability in get_file_format_capabilities(params=params) if capability["enabled"]
    )


def ensure_supported_file_extension(file_name: str | os.PathLike[str], params: dict[str, Any] | None = None) -> None:
    suffix = Path(file_name).suffix.lower()
    capability = get_file_format_capability(suffix, params=params)
    if capability is None:
        raise ValueError(f"Unsupported file type: {suffix or 'no extension'}")
    if not capability["enabled"]:
        message = str(capability["reason"] or f"当前无法处理 {suffix} 文件")
        if suffix in LEGACY_OFFICE_FORMATS:
            code = (
                "converter_unavailable"
                if capability["availability"] == "converter_unavailable"
                else "unsupported_format"
            )
            raise LegacyOfficeConversionError(code, message)
        if suffix in {".gif", ".webp"}:
            raise ImageNormalizationError("unsupported_format", message)
        raise ValueError(message)


@dataclass(slots=True)
class DocumentBlock:
    """A source-aware block produced by a document parser."""

    block_type: str
    order: int
    text: str | None = None
    markdown: str | None = None
    page_number: int | None = None
    sheet_name: str | None = None
    slide_number: int | None = None
    start_char_pos: int | None = None
    end_char_pos: int | None = None
    bbox: list[float] | None = None
    confidence: float | None = None
    table_information: dict[str, Any] | None = None
    image_reference: str | None = None
    parser_name: str | None = None
    parser_version: str | None = None
    warnings: list[str] | None = None

    def to_dict(self) -> dict[str, Any]:
        return {
            key: value
            for key, value in {
                "block_type": self.block_type,
                "order": self.order,
                "text": self.text,
                "markdown": self.markdown,
                "page_number": self.page_number,
                "sheet_name": self.sheet_name,
                "slide_number": self.slide_number,
                "start_char_pos": self.start_char_pos,
                "end_char_pos": self.end_char_pos,
                "bbox": self.bbox,
                "confidence": self.confidence,
                "table_information": self.table_information,
                "image_reference": self.image_reference,
                "parser_name": self.parser_name,
                "parser_version": self.parser_version,
                "warnings": self.warnings,
            }.items()
            if value is not None
        }


@dataclass(slots=True)
class MarkdownParseResult:
    """Unified Markdown and source-structure parsing result."""

    markdown: str
    document_title: str | None
    parser_name: str
    parser_version: str
    warnings: list[str] = field(default_factory=list)
    blocks: list[DocumentBlock] = field(default_factory=list)
    file_ext: str | None = None
    artifacts: dict[str, Any] = field(default_factory=dict)
    classification: dict[str, Any] = field(default_factory=dict)
    attempts: list[dict[str, Any]] = field(default_factory=list)
    quality: dict[str, Any] = field(default_factory=dict)
    format_metadata: dict[str, Any] = field(default_factory=dict)

    def to_metadata(self) -> dict[str, Any]:
        """Serialize metadata used later to map chunks back to source blocks."""
        metadata = {
            "document_title": self.document_title,
            "parser_name": self.parser_name,
            "parser_version": self.parser_version,
            "warnings": list(self.warnings),
            "blocks": [block.to_dict() for block in self.blocks],
            "file_ext": self.file_ext,
            "classification": dict(self.classification),
            "attempts": list(self.attempts),
            "quality": dict(self.quality),
        }
        metadata.update(self.format_metadata)
        return metadata


ParseResult = MarkdownParseResult

_docling_converter: DocumentConverter | None = None
_docling_converter_lock = threading.Lock()


def _package_version(package_name: str) -> str:
    try:
        return version(package_name)
    except PackageNotFoundError:
        return "unknown"


def _contains_semantic_text(value: str) -> bool:
    return bool(re.search(r"[A-Za-z0-9\u3400-\u9fff]", value or ""))


def _minimum_valid_characters(params: dict | None = None) -> int:
    configured = (params or {}).get("min_valid_text_chars")
    if configured is None:
        configured = os.getenv("DOCUMENT_PARSE_MIN_VALID_CHARS", "1")
    try:
        return max(int(configured), 1)
    except (TypeError, ValueError):
        return 1


def validate_parse_result(result: MarkdownParseResult, params: dict | None = None) -> MarkdownParseResult:
    """Reject empty or placeholder-only output before Markdown is persisted."""
    semantic_chars = re.findall(r"[A-Za-z0-9\u3400-\u9fff]", result.markdown or "")
    valid_blocks = [block for block in result.blocks if _contains_semantic_text(block.text or block.markdown or "")]
    if len(semantic_chars) < _minimum_valid_characters(params) or not valid_blocks:
        raise ValueError("文档未提取到有效文本，解析结果为空")
    return result


def _blocks_from_markdown(markdown: str) -> tuple[str | None, list[DocumentBlock]]:
    blocks: list[DocumentBlock] = []
    title: str | None = None
    pattern = r"(?:^|\n\s*\n)(?P<block>\s*\S[\s\S]*?)(?=\n\s*\n|\Z)"
    for match in re.finditer(pattern, markdown):
        raw = match.group("block")
        content = raw.strip()
        if not content:
            continue
        heading = re.match(r"^(#{1,6})\s+(.+)$", content.splitlines()[0])
        is_table = "|" in content and any(re.match(r"^\s*\|?\s*:?-{3,}", line) for line in content.splitlines()[1:2])
        block_type = "table" if is_table else "heading" if heading else "paragraph"
        text = heading.group(2).strip() if heading else content
        if heading and title is None:
            title = text
        start = match.start("block") + len(raw) - len(raw.lstrip())
        blocks.append(
            DocumentBlock(
                block_type=block_type,
                order=len(blocks),
                text=text,
                markdown=content,
                start_char_pos=start,
                end_char_pos=start + len(content),
            )
        )
    return title, blocks


def _markdown_table(rows: list[list[Any]]) -> str:
    width = max((len(row) for row in rows), default=0)
    normalized = [
        [str(value if value is not None else "").replace("|", r"\|").replace("\n", " ").strip() for value in row]
        + [""] * (width - len(row))
        for row in rows
    ]
    while normalized and not any(normalized[-1]):
        normalized.pop()
    if not normalized or not any(any(row) for row in normalized):
        return ""
    header = normalized[0]
    lines = [f"| {' | '.join(header)} |", f"| {' | '.join(['---'] * width)} |"]
    lines.extend(f"| {' | '.join(row)} |" for row in normalized[1:])
    return "\n".join(lines)


def _compose_blocks(
    blocks: list[DocumentBlock],
    prefixes: dict[int, str] | None = None,
) -> str:
    rendered: list[str] = []
    cursor = 0
    prefixes = prefixes or {}
    for block in blocks:
        prefix = prefixes.get(block.order, "")
        value = block.markdown or block.text or ""
        part = f"{prefix}{value}"
        if rendered:
            cursor += 2
        block.start_char_pos = cursor + len(prefix)
        block.end_char_pos = block.start_char_pos + len(value)
        rendered.append(part)
        cursor += len(part)
    return "\n\n".join(rendered)


def _validate_pdf_bytes(content: bytes, policy: OCRRoutingPolicy) -> None:
    if not content.startswith(b"%PDF-"):
        raise ValueError("PDF 文件签名不匹配")
    try:
        document = fitz.open(stream=content, filetype="pdf")
        try:
            if document.needs_pass:
                raise ValueError("PDF 文件已加密，无法解析")
            if document.page_count <= 0:
                raise ValueError("PDF 文件不包含可解析页面")
            if document.page_count > policy.max_pdf_pages:
                raise ValueError(f"PDF 页数超过安全限制（最多 {policy.max_pdf_pages} 页）")
            document.load_page(0)
        finally:
            document.close()
    except ValueError:
        raise
    except Exception as exc:
        raise ValueError("PDF 文件损坏或无法读取") from exc


def _validate_image_bytes(suffix: str, content: bytes, policy: OCRRoutingPolicy) -> None:
    validate_image_bytes(suffix, content, policy)


def validate_document_bytes(filename: str, content: bytes, params: dict | None = None) -> None:
    """Validate runtime support and basic file/container signatures."""
    policy = OCRRoutingPolicy.from_params(params)
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_FILE_EXTENSIONS:
        raise ValueError(f"不支持的文件格式: {suffix or '无扩展名'}")
    if suffix in LEGACY_OFFICE_FORMATS:
        validate_legacy_office_bytes(suffix, content, params=params)
        return
    ensure_supported_file_extension(filename, params=params)
    if suffix in {".txt", ".md"}:
        if b"\x00" in content:
            raise ValueError("文本文件包含二进制内容")
        try:
            content.decode("utf-8-sig")
        except UnicodeDecodeError as exc:
            raise ValueError("文本文件必须使用 UTF-8 编码") from exc
        return
    if suffix == ".pdf":
        _validate_pdf_bytes(content, policy)
        return
    if suffix in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
        _validate_image_bytes(suffix, content, policy)
        return
    validate_ooxml_bytes(suffix, content)


def _get_docling_converter() -> DocumentConverter:
    global _docling_converter
    if _docling_converter is None:
        _docling_converter = DocumentConverter(
            format_options={
                InputFormat.DOCX: None,
                InputFormat.XLSX: None,
                InputFormat.PPTX: None,
            }
        )
    return _docling_converter


def _resolve_image_storage_params(params: dict | None) -> tuple[str, str]:
    params = params or {}
    image_bucket = params.get("image_bucket") or "public"
    image_prefix = params.get("image_prefix")
    if image_prefix:
        normalized_prefix = str(image_prefix).strip("/")
        if normalized_prefix:
            return image_bucket, normalized_prefix
    return image_bucket, "unknown/kb-images"


def _resolve_ocr_engine_params(params: dict | None) -> tuple[str, dict[str, Any]]:
    from yuxi import config

    params = params or {}
    engine = str(params.get("ocr_engine") if "ocr_engine" in params else config.default_ocr_engine)
    engine = engine.strip() or config.default_ocr_engine
    engine_config = params.get("ocr_engine_config")
    processor_params = dict(params)
    if isinstance(engine_config, dict):
        processor_params.update(engine_config)
    return engine, processor_params


def _upload_image_to_minio(image_data: bytes, filename: str, bucket_name: str, object_prefix: str) -> str:
    minio_client = get_minio_client()
    minio_client.ensure_bucket_exists(bucket_name)
    normalized_prefix = object_prefix.strip("/") or "unknown/kb-images"
    timestamp = int(time.time() * 1000000)
    object_name = f"{normalized_prefix}/{timestamp}_{Path(filename).name}"
    result = minio_client.upload_file(
        bucket_name=bucket_name,
        object_name=object_name,
        data=image_data,
    )
    return result.url


def _parse_data_uri(data_uri: str) -> tuple[bytes, str]:
    header, base64_data = data_uri.split(",", 1)
    mime_type = header.split(":")[1].split(";")[0]
    return base64.b64decode(base64_data), mime_type


def _convert_with_docling(file_path: Path, params: dict | None = None) -> str:
    """Convert Office Open XML through the existing Docling adapter."""
    params = params or {}
    image_bucket, image_prefix = _resolve_image_storage_params(params)
    with _docling_converter_lock:
        result = _get_docling_converter().convert(file_path)
    if result.status.name != "SUCCESS":
        raise RuntimeError(f"Docling 转换失败: {result.status}")
    doc = result.document
    if not getattr(doc, "pictures", None):
        return doc.export_to_markdown()

    replacements: list[str] = []
    for picture in doc.pictures:
        uri = str(picture.image.uri) if hasattr(picture, "image") and hasattr(picture.image, "uri") else ""
        if not uri.startswith("data:"):
            replacements.append("")
            continue
        filename = "image"
        try:
            image_data, mime_type = _parse_data_uri(uri)
            filename = f"image_{int(time.time() * 1000000)}.{mime_type.split('/')[-1]}"
            url = _upload_image_to_minio(image_data, filename, image_bucket, image_prefix)
            replacements.append(f"![{filename}]({url})")
        except Exception as exc:  # noqa: BLE001
            logger.error("上传解析图片失败 {}: {}", filename, exc)
            replacements.append(f"[图片: {filename}]")

    markdown = doc.export_to_markdown()
    for replacement in replacements:
        markdown = re.sub(r"<!--\s*image\s*-->", replacement, markdown, count=1)
    return markdown


def _parse_docx_fallback(file_path: Path) -> MarkdownParseResult:
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(str(file_path))
    blocks: list[DocumentBlock] = []
    title: str | None = None
    for child in document.element.body.iterchildren():
        if child.tag.endswith("}p"):
            paragraph = Paragraph(child, document)
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = str(getattr(paragraph.style, "name", "") or "")
            heading_match = re.match(r"Heading\s+(\d+)", style_name, re.IGNORECASE)
            block_type = "heading" if heading_match or style_name.lower() == "title" else "paragraph"
            level = int(heading_match.group(1)) if heading_match else 1
            markdown = f"{'#' * min(max(level, 1), 6)} {text}" if block_type == "heading" else text
            if block_type == "heading" and title is None:
                title = text
            blocks.append(
                DocumentBlock(
                    block_type=block_type,
                    order=len(blocks),
                    text=text,
                    markdown=markdown,
                )
            )
        elif child.tag.endswith("}tbl"):
            table = Table(child, document)
            markdown = _markdown_table([[cell.text for cell in row.cells] for row in table.rows])
            if markdown:
                blocks.append(
                    DocumentBlock(
                        block_type="table",
                        order=len(blocks),
                        text="\n".join(cell.text for row in table.rows for cell in row.cells),
                        markdown=markdown,
                    )
                )
    markdown = _compose_blocks(blocks)
    return MarkdownParseResult(
        markdown=markdown,
        document_title=title or file_path.stem,
        parser_name="python-docx",
        parser_version=_package_version("python-docx"),
        warnings=["Docling 解析失败，已回退到 python-docx"],
        blocks=blocks,
        file_ext=".docx",
    )


def _parse_docx(file_path: Path, params: dict | None = None) -> MarkdownParseResult:
    try:
        markdown = _convert_with_docling(file_path, params=params)
    except Exception as exc:  # noqa: BLE001
        logger.warning("Docling 解析 DOCX 失败，回退到 python-docx: {}, {}", file_path.name, exc)
        return _parse_docx_fallback(file_path)
    title, blocks = _blocks_from_markdown(markdown)
    return MarkdownParseResult(
        markdown=markdown,
        document_title=title or file_path.stem,
        parser_name="docling",
        parser_version=_package_version("docling"),
        blocks=blocks,
        file_ext=".docx",
    )


def _convert_docx_with_python_docx(file_path: Path) -> str:
    """Compatibility helper retained for existing callers and tests."""
    return _parse_docx_fallback(file_path).markdown


def _parse_xlsx(file_path: Path) -> MarkdownParseResult:
    from openpyxl import load_workbook

    workbook = load_workbook(file_path, read_only=True, data_only=True)
    blocks: list[DocumentBlock] = []
    prefixes: dict[int, str] = {}
    try:
        for sheet in workbook.worksheets:
            rows = [list(row) for row in sheet.iter_rows(values_only=True)]
            markdown = _markdown_table(rows)
            if not markdown:
                continue
            order = len(blocks)
            prefixes[order] = f"## {sheet.title}\n\n"
            blocks.append(
                DocumentBlock(
                    block_type="table",
                    order=order,
                    text="\n".join(
                        str(value) for row in rows for value in row if value is not None and str(value).strip()
                    ),
                    markdown=markdown,
                    sheet_name=sheet.title,
                )
            )
    finally:
        workbook.close()
    return MarkdownParseResult(
        markdown=_compose_blocks(blocks, prefixes),
        document_title=file_path.stem,
        parser_name="openpyxl",
        parser_version=_package_version("openpyxl"),
        blocks=blocks,
        file_ext=".xlsx",
    )


def _parse_pptx(file_path: Path) -> MarkdownParseResult:
    from pptx import Presentation

    presentation = Presentation(str(file_path))
    blocks: list[DocumentBlock] = []
    prefixes: dict[int, str] = {}
    title: str | None = None
    for slide_number, slide in enumerate(presentation.slides, 1):
        slide_start = len(blocks)
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                markdown = _markdown_table([[cell.text for cell in row.cells] for row in shape.table.rows])
                if markdown:
                    blocks.append(
                        DocumentBlock(
                            block_type="table",
                            order=len(blocks),
                            text="\n".join(cell.text for row in shape.table.rows for cell in row.cells),
                            markdown=markdown,
                            slide_number=slide_number,
                        )
                    )
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if not text:
                continue
            is_title = shape == slide.shapes.title
            block_type = "title" if is_title else "paragraph"
            markdown = f"### {text}" if is_title else text
            if is_title and title is None:
                title = text
            blocks.append(
                DocumentBlock(
                    block_type=block_type,
                    order=len(blocks),
                    text=text,
                    markdown=markdown,
                    slide_number=slide_number,
                )
            )
        if len(blocks) > slide_start:
            prefixes[slide_start] = f"## Slide {slide_number}\n\n"
    return MarkdownParseResult(
        markdown=_compose_blocks(blocks, prefixes),
        document_title=title or file_path.stem,
        parser_name="python-pptx",
        parser_version=_package_version("python-pptx"),
        blocks=blocks,
        file_ext=".pptx",
    )


def _parse_text_pdf(file_path: Path) -> MarkdownParseResult:
    docs = PyPDFLoader(str(file_path)).load()
    blocks: list[DocumentBlock] = []
    prefixes: dict[int, str] = {}
    for default_index, doc in enumerate(docs):
        text = (doc.page_content or "").strip()
        if not text:
            continue
        page_number = int((doc.metadata or {}).get("page", default_index)) + 1
        order = len(blocks)
        prefixes[order] = f"<!-- page: {page_number} -->\n\n"
        blocks.append(
            DocumentBlock(
                block_type="paragraph",
                order=order,
                text=text,
                markdown=text,
                page_number=page_number,
            )
        )
    if not blocks:
        raise ValueError("文本型 PDF 未提取到有效文本，可能是扫描 PDF，需要启用 OCR")
    return MarkdownParseResult(
        markdown=_compose_blocks(blocks, prefixes),
        document_title=file_path.stem,
        parser_name="native_pdf",
        parser_version=_package_version("pypdf"),
        blocks=blocks,
        file_ext=".pdf",
    )


def _page_text_coverage(page: fitz.Page) -> float:
    page_area = max(float(page.rect.width * page.rect.height), 1.0)
    covered_area = 0.0
    for block in page.get_text("blocks"):
        if len(block) < 5 or not _contains_semantic_text(str(block[4])):
            continue
        rectangle = fitz.Rect(block[:4]) & page.rect
        if not rectangle.is_empty:
            covered_area += max(float(rectangle.width * rectangle.height), 0.0)
    return min(covered_area / page_area, 1.0)


def _render_pdf_page_for_ocr(
    page: fitz.Page,
    *,
    page_number: int,
    policy: OCRRoutingPolicy,
) -> Path:
    scale = policy.pdf_render_scale
    width = max(int(page.rect.width * scale), 1)
    height = max(int(page.rect.height * scale), 1)
    if width * height > policy.max_pdf_page_pixels:
        raise ValueError(f"PDF 第 {page_number} 页渲染像素超过安全限制（最多 {policy.max_pdf_page_pixels}）")
    pixmap = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temporary:
        output_path = Path(temporary.name)
    try:
        pixmap.save(output_path)
        return output_path
    except Exception:
        output_path.unlink(missing_ok=True)
        raise


def _ocr_route_block(route: OCRRouteResult, *, order: int, page_number: int) -> DocumentBlock:
    quality = route.quality or {}
    return DocumentBlock(
        block_type="table" if int(quality.get("table_valid_cells") or 0) > 0 else "paragraph",
        order=order,
        text=route.markdown,
        markdown=route.markdown,
        page_number=page_number,
        parser_name=route.parser_name,
        parser_version=route.parser_version,
        warnings=list(route.warnings),
    )


async def _parse_pdf_auto(file_path: Path, params: dict | None = None) -> MarkdownParseResult:
    options = dict(params or {})
    policy = OCRRoutingPolicy.from_params(options)
    options["_ocr_deadline_monotonic"] = time.monotonic() + policy.max_ocr_seconds
    await notify_processing_stage(options, "detecting", 15)
    document = fitz.open(file_path)
    blocks: list[DocumentBlock] = []
    prefixes: dict[int, str] = {}
    attempts: list[dict[str, Any]] = []
    page_classifications: list[dict[str, Any]] = []
    warnings_list: list[str] = []
    native_pages = 0
    ocr_pages = 0
    ocr_parser_names: set[str] = set()
    try:
        if document.page_count > policy.max_pdf_pages:
            raise ValueError(f"PDF 页数超过安全限制（最多 {policy.max_pdf_pages} 页）")
        for page_index in range(document.page_count):
            page_number = page_index + 1
            page = document.load_page(page_index)
            await notify_processing_stage(options, "extracting_text", 25)
            native_started = time.monotonic()
            native_text = str(page.get_text("text") or "").strip()
            coverage = _page_text_coverage(page)
            native_quality = assess_text_quality(
                native_text,
                policy=policy,
                min_valid_characters=policy.native_pdf_min_valid_characters,
                page_coverage=coverage,
            )
            native_attempt = {
                "provider": "native_pdf",
                "stage": "extracting_text",
                "status": "accepted" if native_quality.accepted else "rejected",
                "duration_ms": round((time.monotonic() - native_started) * 1000),
                "page_number": page_number,
                "quality": native_quality.to_dict(),
            }
            attempts.append(native_attempt)
            page_classification = {
                "page_number": page_number,
                "native_quality": native_quality.to_dict(),
                "route": "native_text" if native_quality.accepted else "ocr",
                "reasons": native_quality.reasons or ["native_text_quality_accepted"],
            }
            page_classifications.append(page_classification)

            if native_quality.accepted:
                order = len(blocks)
                prefixes[order] = f"<!-- page: {page_number} -->\n\n"
                blocks.append(
                    DocumentBlock(
                        block_type="paragraph",
                        order=order,
                        text=native_text,
                        markdown=native_text,
                        page_number=page_number,
                        parser_name="native_pdf",
                        parser_version=_package_version("PyMuPDF"),
                        warnings=[],
                    )
                )
                native_pages += 1
                continue

            has_visual_content = bool(page.get_images(full=True) or page.get_drawings())
            if not has_visual_content:
                native_attempt["status"] = "skipped"
                native_attempt["failure_reason"] = "blank_page_without_visual_content"
                page_classification["route"] = "blank_ignored"
                page_classification["reasons"] = ["blank_page_without_visual_content"]
                continue

            rendered_page = _render_pdf_page_for_ocr(
                page,
                page_number=page_number,
                policy=policy,
            )
            try:
                page_options = dict(options)
                page_options["_complex_layout"] = looks_like_structured_layout(rendered_page)
                try:
                    route = await run_ocr_fallback(
                        rendered_page,
                        params=page_options,
                        page_number=page_number,
                    )
                except OCRRoutingError as exc:
                    combined_attempts = [
                        *attempts,
                        *[
                            {**attempt, "page_number": page_number}
                            for attempt in exc.parse_metadata.get("attempts", [])
                        ],
                    ]
                    error = OCRRoutingError(
                        f"PDF 第 {page_number} 页 OCR 未提取到有效文本",
                        attempts=combined_attempts,
                        warnings=[*warnings_list, *exc.parse_metadata.get("warnings", [])],
                    )
                    error.parse_metadata["classification"] = {
                        "type": "pdf_ocr_failed",
                        "pages": page_classifications,
                        "thresholds": policy.public_metadata(),
                    }
                    raise error from exc
            finally:
                rendered_page.unlink(missing_ok=True)

            order = len(blocks)
            prefixes[order] = f"<!-- page: {page_number} -->\n\n"
            blocks.append(_ocr_route_block(route, order=order, page_number=page_number))
            attempts.extend({**attempt, "page_number": page_number} for attempt in route.attempts)
            warnings_list.extend(route.warnings)
            ocr_parser_names.add(route.parser_name)
            ocr_pages += 1
    finally:
        document.close()

    if native_pages and ocr_pages:
        classification_type = "mixed_pdf"
        parser_name = "hybrid_pdf"
    elif ocr_pages:
        classification_type = "scanned_pdf"
        parser_name = next(iter(ocr_parser_names)) if len(ocr_parser_names) == 1 else "hybrid_pdf"
    else:
        classification_type = "text_pdf"
        parser_name = "native_pdf"
    markdown = _compose_blocks(blocks, prefixes)
    quality = assess_text_quality(markdown, policy=policy)
    return MarkdownParseResult(
        markdown=markdown,
        document_title=file_path.stem,
        parser_name=parser_name,
        parser_version=_package_version("PyMuPDF") if parser_name == "native_pdf" else "1",
        warnings=list(dict.fromkeys(warnings_list)),
        blocks=blocks,
        file_ext=".pdf",
        classification={
            "type": classification_type,
            "reason": (
                "all_pages_passed_native_text_quality"
                if classification_type == "text_pdf"
                else "all_pages_required_ocr"
                if classification_type == "scanned_pdf"
                else "native_and_ocr_pages_combined"
            ),
            "native_page_count": native_pages,
            "ocr_page_count": ocr_pages,
            "empty_page_ratio": round(
                sum(not page["native_quality"]["valid_characters"] for page in page_classifications)
                / max(len(page_classifications), 1),
                4,
            ),
            "pages": page_classifications,
            "thresholds": policy.public_metadata(),
        },
        attempts=attempts,
        quality=quality.to_dict(),
    )


async def _parse_image_auto(file_path: Path, params: dict | None = None) -> MarkdownParseResult:
    options = dict(params or {})
    policy = OCRRoutingPolicy.from_params(options)
    await notify_processing_stage(options, "detecting", 15)
    normalized_image = None
    ocr_path = file_path
    format_metadata: dict[str, Any] = {}
    try:
        if file_path.suffix.lower() in {".gif", ".webp"}:
            normalized_image = await asyncio.to_thread(normalize_image_for_ocr, file_path, options)
            ocr_path = normalized_image.path
            format_metadata = dict(normalized_image.metadata)
        options["_complex_layout"] = looks_like_structured_layout(ocr_path)
        try:
            route = await run_ocr_fallback(ocr_path, params=options, page_number=1)
        except OCRRoutingError as exc:
            exc.parse_metadata["classification"] = {
                "type": "image",
                "reason": "ocr_quality_not_accepted",
                "thresholds": policy.public_metadata(),
            }
            exc.parse_metadata.update(format_metadata)
            raise
        block = _ocr_route_block(route, order=0, page_number=1)
        markdown = _compose_blocks([block])
        return MarkdownParseResult(
            markdown=markdown,
            document_title=file_path.stem,
            parser_name=route.parser_name,
            parser_version=route.parser_version,
            warnings=route.warnings,
            blocks=[block],
            file_ext=file_path.suffix.lower(),
            classification={
                "type": "image",
                "reason": "automatic_ocr_route",
                "structured_layout_detected": bool(options["_complex_layout"]),
                "thresholds": policy.public_metadata(),
            },
            attempts=route.attempts,
            quality=route.quality,
            format_metadata=format_metadata,
        )
    finally:
        if normalized_image is not None:
            await asyncio.to_thread(normalized_image.cleanup)


async def _parse_legacy_office(
    file_path: Path,
    params: dict | None = None,
) -> MarkdownParseResult:
    options = dict(params or {})
    await notify_processing_stage(options, "converting", 25)
    conversion = await asyncio.to_thread(LegacyOfficeConverter(options).convert, file_path)
    normalized_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=conversion.normalized_suffix) as temp_file:
            normalized_path = Path(temp_file.name)
            temp_file.write(conversion.content)
        await notify_processing_stage(options, "extracting_text", 35)
        if conversion.normalized_suffix == ".docx":
            result = await asyncio.to_thread(_parse_docx, normalized_path, options)
        elif conversion.normalized_suffix == ".xlsx":
            result = await asyncio.to_thread(_parse_xlsx, normalized_path)
        else:
            result = await asyncio.to_thread(_parse_pptx, normalized_path)
        if result.document_title == normalized_path.stem:
            result.document_title = file_path.stem
        result.file_ext = file_path.suffix.lower()
        result.warnings = [*conversion.metadata.get("conversion_warnings", []), *result.warnings]
        result.format_metadata.update(conversion.metadata)
        return result
    finally:
        if normalized_path is not None:
            normalized_path.unlink(missing_ok=True)


def _parse_text_file(file_path: Path, *, markdown_source: bool) -> MarkdownParseResult:
    try:
        content = file_path.read_bytes().decode("utf-8-sig")
    except UnicodeDecodeError as exc:
        raise ValueError("文本文件必须使用 UTF-8 编码") from exc
    if markdown_source:
        title, blocks = _blocks_from_markdown(content)
        markdown = content
        parser_name = "native_markdown"
    else:
        blocks = []
        cursor = 0
        for paragraph in re.split(r"\n\s*\n", content):
            text = paragraph.strip()
            if not text:
                continue
            start = content.find(paragraph, cursor) + len(paragraph) - len(paragraph.lstrip())
            cursor = start + len(text)
            blocks.append(
                DocumentBlock(
                    block_type="paragraph",
                    order=len(blocks),
                    text=text,
                    markdown=text,
                    start_char_pos=start,
                    end_char_pos=start + len(text),
                )
            )
        title = None
        markdown = content.strip()
        parser_name = "native_text"
    return MarkdownParseResult(
        markdown=markdown,
        document_title=title or file_path.stem,
        parser_name=parser_name,
        parser_version="1",
        blocks=blocks,
        file_ext=file_path.suffix.lower(),
    )


def _convert_csv_to_markdown(file_path: Path) -> str:
    import pandas as pd

    dataframe = pd.read_csv(file_path)
    return "\n\n".join(dataframe.iloc[[index]].to_markdown(index=False) for index in range(len(dataframe)))


def pdfreader(file_path, params=None):
    """Compatibility native PDF text reader."""
    del params
    path = Path(file_path)
    assert path.exists(), "File not found"
    assert path.suffix.lower() == ".pdf", "File format not supported"
    return "\n\n".join(doc.page_content for doc in PyPDFLoader(str(path)).load())


def parse_pdf(file, params=None):
    """Existing configurable OCR entry retained outside the stable text-PDF route."""
    from yuxi.knowledge.parser.base import DocumentProcessorException
    from yuxi.knowledge.parser.factory import DocumentProcessorFactory

    opt_ocr, processor_params = _resolve_ocr_engine_params(params)
    if opt_ocr == "disable":
        return pdfreader(file, params=processor_params)
    image_bucket, image_prefix = _resolve_image_storage_params(processor_params)
    processor_params.setdefault("image_bucket", image_bucket)
    processor_params.setdefault("image_prefix", image_prefix)
    try:
        return DocumentProcessorFactory.process_file(opt_ocr, file, processor_params)
    except DocumentProcessorException:
        raise
    except Exception as exc:  # noqa: BLE001
        raise DocumentProcessorException(f"PDF解析失败: {exc}", opt_ocr, "parsing_failed") from exc


def parse_image(file, params=None):
    """Existing configurable image OCR entry retained for the later OCR PR."""
    from yuxi.knowledge.parser.base import DocumentProcessorException
    from yuxi.knowledge.parser.factory import DocumentProcessorFactory

    opt_ocr, processor_params = _resolve_ocr_engine_params(params)
    if opt_ocr == "disable":
        raise ValueError("图像文件必须启用OCR才能提取文本内容")
    image_bucket, image_prefix = _resolve_image_storage_params(processor_params)
    processor_params.setdefault("image_bucket", image_bucket)
    processor_params.setdefault("image_prefix", image_prefix)
    try:
        return DocumentProcessorFactory.process_file(opt_ocr, file, processor_params)
    except DocumentProcessorException:
        raise
    except Exception as exc:  # noqa: BLE001
        raise DocumentProcessorException(f"图像解析失败: {exc}", opt_ocr, "parsing_failed") from exc


async def parse_pdf_async(file, params=None):
    return await asyncio.to_thread(parse_pdf, file, params=params)


async def parse_image_async(file, params=None):
    return await asyncio.to_thread(parse_image, file, params=params)


async def _download_source_if_needed(file_path: str) -> tuple[Path, bool]:
    from yuxi.knowledge.utils.kb_utils import is_minio_url, parse_minio_url
    from yuxi.storage.minio.client import get_minio_client

    if not is_minio_url(file_path):
        return Path(file_path), False
    clean_path = file_path.split("?", 1)[0]
    suffix = Path(clean_path.rsplit("/", 1)[-1]).suffix
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
        temp_path = Path(temp_file.name)
    try:
        bucket_name, object_name = parse_minio_url(file_path)
        content = await get_minio_client().adownload_file(bucket_name, object_name)
        async with aiofiles.open(temp_path, "wb") as stream:
            await stream.write(content)
        return temp_path, True
    except Exception:
        temp_path.unlink(missing_ok=True)
        raise


async def _process_file_to_result_core(
    file_path: str,
    params: dict | None = None,
) -> MarkdownParseResult:
    actual_path, temporary = await _download_source_if_needed(file_path)
    try:
        suffix = actual_path.suffix.lower()
        ensure_supported_file_extension(actual_path, params=params)
        if suffix == ".txt":
            await notify_processing_stage(params or {}, "extracting_text", 30)
            result = await asyncio.to_thread(_parse_text_file, actual_path, markdown_source=False)
        elif suffix == ".md":
            await notify_processing_stage(params or {}, "extracting_text", 30)
            result = await asyncio.to_thread(_parse_text_file, actual_path, markdown_source=True)
        elif suffix == ".pdf":
            result = await _parse_pdf_auto(actual_path, params)
        elif suffix == ".docx":
            await notify_processing_stage(params or {}, "extracting_text", 30)
            result = await asyncio.to_thread(_parse_docx, actual_path, params)
        elif suffix == ".xlsx":
            await notify_processing_stage(params or {}, "extracting_text", 30)
            result = await asyncio.to_thread(_parse_xlsx, actual_path)
        elif suffix == ".pptx":
            await notify_processing_stage(params or {}, "extracting_text", 30)
            result = await asyncio.to_thread(_parse_pptx, actual_path)
        elif suffix in LEGACY_OFFICE_FORMATS:
            result = await _parse_legacy_office(actual_path, params)
        elif suffix in {".jpg", ".jpeg", ".png", ".gif", ".webp"}:
            result = await _parse_image_auto(actual_path, params)
        elif suffix in {".html", ".htm"}:
            async with aiofiles.open(actual_path, encoding="utf-8") as stream:
                markdown = await asyncio.to_thread(md_convert, await stream.read(), heading_style="ATX")
            title, blocks = _blocks_from_markdown(markdown)
            result = MarkdownParseResult(
                markdown=markdown,
                document_title=title or actual_path.stem,
                parser_name="markdownify",
                parser_version=_package_version("markdownify"),
                blocks=blocks,
                file_ext=suffix,
            )
        elif suffix == ".csv":
            markdown = await asyncio.to_thread(_convert_csv_to_markdown, actual_path)
            title, blocks = _blocks_from_markdown(markdown)
            result = MarkdownParseResult(
                markdown=markdown,
                document_title=actual_path.stem,
                parser_name="pandas",
                parser_version=_package_version("pandas"),
                blocks=blocks,
                file_ext=suffix,
            )
        elif suffix == ".json":
            import json

            async with aiofiles.open(actual_path, encoding="utf-8") as stream:
                data = json.loads(await stream.read())
            markdown = f"```json\n{json.dumps(data, ensure_ascii=False, indent=2)}\n```"
            _, blocks = _blocks_from_markdown(markdown)
            result = MarkdownParseResult(
                markdown=markdown,
                document_title=actual_path.stem,
                parser_name="json",
                parser_version="1",
                blocks=blocks,
                file_ext=suffix,
            )
        elif suffix == ".zip":
            image_bucket, image_prefix = _resolve_image_storage_params(params)
            zip_result = await _process_zip_file(
                str(actual_path),
                image_bucket=image_bucket,
                image_prefix=image_prefix,
            )
            markdown = zip_result["markdown_content"]
            title, blocks = _blocks_from_markdown(markdown)
            result = MarkdownParseResult(
                markdown=markdown,
                document_title=title or actual_path.stem,
                parser_name="zip",
                parser_version="1",
                blocks=blocks,
                file_ext=suffix,
                artifacts={
                    "zip_images_info": zip_result["images_info"],
                    "zip_content_hash": zip_result["content_hash"],
                    "zip_image_bucket": image_bucket,
                    "zip_image_prefix": image_prefix,
                },
            )
        else:
            raise ValueError(f"Unsupported file type: {suffix}")
        return validate_parse_result(result, params)
    finally:
        if temporary:
            actual_path.unlink(missing_ok=True)


async def _process_file_to_markdown_core(
    file_path: str,
    params: dict | None = None,
) -> tuple[str, str | None, dict[str, Any]]:
    """Compatibility adapter for callers that still expect the old tuple."""
    result = await _process_file_to_result_core(file_path, params=params)
    return result.markdown, result.file_ext, result.artifacts


async def parse_source_to_markdown(source: str, params: dict | None = None) -> MarkdownParseResult:
    """Parse a source into the unified result (legacy function name retained)."""
    return await _process_file_to_result_core(source, params=params)


async def parse_source_to_result(source: str, params: dict | None = None) -> MarkdownParseResult:
    return await _process_file_to_result_core(source, params=params)


class Parser:
    """Facade with structured and legacy string-returning parsing methods."""

    @staticmethod
    async def aparse_result(source: str, params: dict | None = None) -> MarkdownParseResult:
        return await parse_source_to_result(source=source, params=params)

    @staticmethod
    async def aparse(source: str, params: dict | None = None) -> str:
        return (await Parser.aparse_result(source=source, params=params)).markdown

    @classmethod
    def parse_result(cls, source: str, params: dict | None = None) -> MarkdownParseResult:
        try:
            asyncio.get_running_loop()
        except RuntimeError:
            return asyncio.run(cls.aparse_result(source=source, params=params))
        raise RuntimeError("当前处于异步上下文，请使用 `await Parser.aparse_result(...)`")

    @classmethod
    def parse(cls, source: str, params: dict | None = None) -> str:
        try:
            asyncio.get_running_loop()
        except RuntimeError:
            return asyncio.run(cls.aparse(source=source, params=params))
        raise RuntimeError("当前处于异步上下文，请使用 `await Parser.aparse(...)`")
