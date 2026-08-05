from __future__ import annotations

import io
import uuid

import fitz
import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from yuxi.knowledge.runtime import knowledge_base
from yuxi.repositories.knowledge_chunk_repository import KnowledgeChunkRepository
from yuxi.repositories.knowledge_file_repository import KnowledgeFileRepository

pytestmark = [pytest.mark.asyncio, pytest.mark.integration]


def _document_bytes(extension: str, marker: str) -> bytes:
    if extension in {".txt", ".md"}:
        value = (
            f"# {marker} heading\n\n{marker} searchable markdown paragraph."
            if extension == ".md"
            else f"{marker} searchable plain text paragraph."
        )
        return value.encode()
    if extension == ".pdf":
        document = fitz.open()
        first = document.new_page()
        first.insert_text((72, 72), f"{marker} searchable PDF page one.")
        document.new_page()
        third = document.new_page()
        third.insert_text((72, 72), f"{marker} PDF page three.")
        value = document.tobytes()
        document.close()
        return value
    if extension == ".docx":
        document = Document()
        document.add_heading(f"{marker} DOCX title", level=1)
        document.add_paragraph(f"{marker} paragraph before table.")
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = marker
        table.cell(1, 1).text = "DOCX table"
        document.add_paragraph(f"{marker} paragraph after table.")
        buffer = io.BytesIO()
        document.save(buffer)
        return buffer.getvalue()
    if extension == ".xlsx":
        workbook = Workbook()
        first = workbook.active
        first.title = "Summary"
        first.append(["Name", "Value"])
        first.append([marker, "XLSX table"])
        detail = workbook.create_sheet("Detail")
        detail.append(["Marker", "Description"])
        detail.append([marker, "second sheet"])
        workbook.create_sheet("Empty")
        buffer = io.BytesIO()
        workbook.save(buffer)
        workbook.close()
        return buffer.getvalue()
    if extension == ".pptx":
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"{marker} PPTX title"
        textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        textbox.text = f"{marker} slide paragraph"
        table = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(4), Inches(1)).table
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = marker
        table.cell(1, 1).text = "PPTX table"
        presentation.slides.add_slide(presentation.slide_layouts[6])
        buffer = io.BytesIO()
        presentation.save(buffer)
        return buffer.getvalue()
    raise AssertionError(f"unsupported fixture extension: {extension}")


async def _upload_and_create(test_client, headers, kb_id: str, filename: str, content: bytes) -> str:
    upload = await test_client.post(
        "/api/knowledge/files/upload",
        params={"kb_id": kb_id},
        files={"file": (filename, content, "application/octet-stream")},
        headers=headers,
    )
    assert upload.status_code == 200, upload.text
    payload = upload.json()
    item = payload["file_path"]
    create = await test_client.post(
        f"/api/knowledge/databases/{kb_id}/documents/add",
        json={
            "items": [item],
            "params": {
                "source_paths": {item: payload["filename"]},
                "content_hashes": {item: "client-value-is-not-trusted"},
                "file_sizes": {item: 1},
            },
        },
        headers=headers,
    )
    assert create.status_code == 200, create.text
    return create.json()["items"][0]["file_id"]


def _result_file_ids(results: list[dict]) -> set[str]:
    return {str(item.get("metadata", {}).get("file_id") or item.get("file_id") or "") for item in results}


async def test_stable_formats_parse_chunk_index_and_retrieve_with_source_metadata(
    test_client,
    admin_headers,
    knowledge_database,
    monkeypatch,
) -> None:
    kb_id = knowledge_database["kb_id"]
    kb = await knowledge_base.aget_kb(kb_id)
    if kb_id not in kb.databases_meta:
        await kb._load_metadata()
    collection = await kb._get_milvus_collection(kb_id)
    assert collection is not None
    embedding_field = next(field for field in collection.schema.fields if field.name == "embedding")
    dimension = int(embedding_field.params["dim"])
    extensions = (".txt", ".md", ".pdf", ".docx", ".xlsx", ".pptx")
    markers = {extension: f"fmt{index}{uuid.uuid4().hex}" for index, extension in enumerate(extensions)}

    def deterministic_vectors(texts: list[str]) -> list[list[float]]:
        vectors: list[list[float]] = []
        for text in texts:
            vector = [0.0] * dimension
            matched = False
            for index, marker in enumerate(markers.values()):
                if marker in text:
                    vector[index] = 1.0
                    matched = True
                    break
            if not matched:
                vector[-1] = 1.0
            vectors.append(vector)
        return vectors

    async def async_embed(texts: list[str]) -> list[list[float]]:
        return deterministic_vectors(texts)

    monkeypatch.setattr(
        kb,
        "_get_embedding_function",
        lambda _model_spec, *, sync=False: deterministic_vectors if sync else async_embed,
    )

    file_repository = KnowledgeFileRepository()
    chunk_repository = KnowledgeChunkRepository()
    for extension in extensions:
        marker = markers[extension]
        filename = f"{marker}{extension}"
        file_id = await _upload_and_create(
            test_client,
            admin_headers,
            kb_id,
            filename,
            _document_bytes(extension, marker),
        )

        parsed = await kb.parse_file(kb_id, file_id)
        assert parsed["status"] == "parsed"
        assert parsed["markdown_file"]
        assert parsed["parse_metadata"]["parser_name"]
        assert parsed["parse_metadata"]["parser_version"]
        assert parsed["parse_metadata"]["blocks"]

        indexed = await kb.index_file(kb_id, file_id)
        assert indexed["status"] == "indexed"
        record = await file_repository.get_by_file_id(file_id)
        assert record is not None
        assert record.chunk_count > 0
        assert record.token_count > 0
        assert record.markdown_file

        chunks = await chunk_repository.list_by_file_id(file_id)
        assert chunks
        assert all(chunk.source_metadata for chunk in chunks)
        assert all(chunk.source_metadata["parser_name"] for chunk in chunks)
        if extension == ".pdf":
            assert {block["page_number"] for block in record.parse_metadata["blocks"]} == {1, 3}
            assert any(chunk.source_metadata.get("page_number") in {1, 3} for chunk in chunks)
        if extension == ".xlsx":
            assert {block["sheet_name"] for block in record.parse_metadata["blocks"]} == {"Summary", "Detail"}
            assert any(chunk.source_metadata.get("sheet_name") for chunk in chunks)
        if extension == ".pptx":
            assert {block["slide_number"] for block in record.parse_metadata["blocks"]} == {1}
            assert any(chunk.source_metadata.get("slide_number") == 1 for chunk in chunks)

        results = await kb.aquery(
            marker,
            kb_id,
            search_mode="vector",
            final_top_k=20,
            similarity_threshold=-1.0,
            use_reranker=False,
        )
        assert file_id in _result_file_ids(results)
        matching = next(item for item in results if file_id in _result_file_ids([item]))
        assert matching["metadata"]["source_metadata"]["parser_name"]


@pytest.mark.parametrize(
    ("filename", "content"),
    [
        ("fake.pdf", b"not-pdf"),
        ("fake.docx", b"not-docx"),
        ("legacy.doc", b"legacy"),
        ("image.png", b"image"),
    ],
)
async def test_upload_rejects_unsupported_or_signature_mismatched_files(
    test_client,
    admin_headers,
    knowledge_database,
    filename: str,
    content: bytes,
) -> None:
    response = await test_client.post(
        "/api/knowledge/files/upload",
        params={"kb_id": knowledge_database["kb_id"]},
        files={"file": (filename, content, "application/octet-stream")},
        headers=admin_headers,
    )

    assert response.status_code == 400
    assert "Traceback" not in response.text
