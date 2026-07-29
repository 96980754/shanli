from __future__ import annotations

import asyncio
import io
import subprocess
import uuid
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches
from yuxi.knowledge.parser.legacy_office import (
    OLE_COMPOUND_FILE_SIGNATURE,
    LegacyOfficeConverter,
)
from yuxi.knowledge.runtime import knowledge_base
from yuxi.knowledge.utils.kb_utils import parse_minio_url
from yuxi.repositories.knowledge_chunk_repository import KnowledgeChunkRepository
from yuxi.repositories.knowledge_file_repository import KnowledgeFileRepository
from yuxi.storage.minio.client import get_minio_client

pytestmark = [pytest.mark.asyncio, pytest.mark.integration]

_LEGACY_FILTERS = {
    ".doc": "doc:MS Word 97",
    ".xls": "xls:MS Excel 97",
    ".ppt": "ppt:MS PowerPoint 97",
}


def _font(size: int = 62):
    for name in ("NotoSansCJK-Regular.ttc", "DejaVuSans.ttf"):
        try:
            return ImageFont.truetype(name, size=size)
        except OSError:
            continue
    return ImageFont.load_default()


def _build_ooxml(path: Path, marker: str) -> None:
    if path.suffix == ".docx":
        document = Document()
        document.add_heading(f"{marker} legacy Word title", level=1)
        document.add_paragraph(f"{marker} legacy Word searchable paragraph.")
        document.save(path)
        return
    if path.suffix == ".xlsx":
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "LegacySheet"
        sheet.append(["Marker", "Description"])
        sheet.append([marker, f"{marker} legacy spreadsheet cell"])
        workbook.save(path)
        workbook.close()
        return
    if path.suffix == ".pptx":
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"{marker} legacy slide"
        textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
        textbox.text = f"{marker} legacy presentation paragraph"
        presentation.save(path)
        return
    raise AssertionError(f"unsupported fixture extension: {path.suffix}")


def _legacy_fixture(tmp_path: Path, legacy_suffix: str, marker: str) -> bytes:
    binary = LegacyOfficeConverter.resolve_binary("soffice")
    assert binary, "真实旧 Office 集成测试要求 API 镜像安装 LibreOffice"
    normalized_suffix = {".doc": ".docx", ".xls": ".xlsx", ".ppt": ".pptx"}[legacy_suffix]
    fixture_dir = tmp_path / f"fixture-{legacy_suffix.removeprefix('.')}-{uuid.uuid4().hex}"
    output_dir = fixture_dir / "output"
    profile_dir = fixture_dir / "profile"
    output_dir.mkdir(parents=True)
    profile_dir.mkdir()
    source = fixture_dir / f"source{normalized_suffix}"
    _build_ooxml(source, marker)
    command = [
        binary,
        "--headless",
        "--nologo",
        "--nofirststartwizard",
        "--nodefault",
        "--nolockcheck",
        f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
        "--convert-to",
        _LEGACY_FILTERS[legacy_suffix],
        "--outdir",
        str(output_dir),
        str(source),
    ]
    completed = subprocess.run(command, capture_output=True, timeout=90, check=False, shell=False)
    assert completed.returncode == 0, f"LibreOffice fixture conversion failed with exit {completed.returncode}"
    target = output_dir / f"source{legacy_suffix}"
    assert target.is_file(), f"LibreOffice did not produce {legacy_suffix} fixture"
    content = target.read_bytes()
    assert content.startswith(OLE_COMPOUND_FILE_SIGNATURE)
    assert len(content) > len(OLE_COMPOUND_FILE_SIGNATURE)
    return content


def _image_bytes(text: str, *, image_format: str, animated: bool = False) -> bytes:
    first = Image.new("RGB", (1800, 420), "white")
    draw = ImageDraw.Draw(first)
    draw.text((80, 145), text, font=_font(), fill="black")
    buffer = io.BytesIO()
    if animated:
        second = Image.new("RGB", first.size, "white")
        second_draw = ImageDraw.Draw(second)
        second_draw.text((80, 145), "IGNORED SECOND FRAME", font=_font(), fill="black")
        first.save(
            buffer,
            format=image_format,
            save_all=True,
            append_images=[second],
            duration=100,
            loop=0,
            quality=95,
        )
    else:
        first.save(buffer, format=image_format, quality=95)
    return buffer.getvalue()


async def _upload_and_create(test_client, headers, kb_id: str, filename: str, content: bytes) -> str:
    upload = await test_client.post(
        "/api/knowledge/files/upload",
        params={"kb_id": kb_id},
        files={"file": (filename, content, "application/octet-stream")},
        headers=headers,
    )
    assert upload.status_code == 200, upload.text
    payload = upload.json()
    item = payload["file_path"]
    create = await test_client.post(
        f"/api/knowledge/databases/{kb_id}/documents/add",
        json={
            "items": [item],
            "params": {
                "source_paths": {item: payload["filename"]},
                "content_hashes": {item: "client-content-hash-is-not-trusted"},
                "file_sizes": {item: 1},
            },
        },
        headers=headers,
    )
    assert create.status_code == 200, create.text
    return create.json()["items"][0]["file_id"]


def _result_file_ids(results: list[dict]) -> set[str]:
    return {str(item.get("metadata", {}).get("file_id") or item.get("file_id") or "") for item in results}


async def _configure_deterministic_embedding(kb, kb_id: str, markers: list[str], monkeypatch) -> None:
    if kb_id not in kb.databases_meta:
        await kb._load_metadata()
    collection = await kb._get_milvus_collection(kb_id)
    assert collection is not None
    embedding_field = next(field for field in collection.schema.fields if field.name == "embedding")
    dimension = int(embedding_field.params["dim"])

    def deterministic_vectors(texts: list[str]) -> list[list[float]]:
        vectors: list[list[float]] = []
        for text in texts:
            normalized = text.upper()
            vector = [0.0] * dimension
            for index, marker in enumerate(markers):
                if marker in normalized:
                    vector[index] = 1.0
                    break
            else:
                vector[-1] = 1.0
            vectors.append(vector)
        return vectors

    async def async_embed(texts: list[str]) -> list[list[float]]:
        return deterministic_vectors(texts)

    monkeypatch.setattr(
        kb,
        "_get_embedding_function",
        lambda _model_spec, *, sync=False: deterministic_vectors if sync else async_embed,
    )


async def test_legacy_office_real_conversion_parse_chunk_index_and_retrieve(
    test_client,
    admin_headers,
    knowledge_database,
    tmp_path: Path,
    monkeypatch,
) -> None:
    kb_id = knowledge_database["kb_id"]
    markers = ["LEGACYALPHA", "LEGACYBRAVO", "LEGACYCHARLIE"]
    formats = [
        (".doc", markers[0]),
        (".xls", markers[1]),
        (".ppt", markers[2]),
    ]
    capabilities = await test_client.get("/api/knowledge/files/supported-types", headers=admin_headers)
    assert capabilities.status_code == 200, capabilities.text
    capability_payload = capabilities.json()
    assert {".doc", ".xls", ".ppt"} <= set(capability_payload["file_types"])
    capability_map = {item["extension"]: item for item in capability_payload["capabilities"]}
    assert all(capability_map[suffix]["availability"] == "available" for suffix, _marker in formats)

    kb = await knowledge_base.aget_kb(kb_id)
    await _configure_deterministic_embedding(kb, kb_id, markers, monkeypatch)
    file_repository = KnowledgeFileRepository()
    chunk_repository = KnowledgeChunkRepository()
    file_ids: list[str] = []
    for suffix, marker in formats:
        content = await asyncio.to_thread(_legacy_fixture, tmp_path, suffix, marker)
        file_id = await _upload_and_create(test_client, admin_headers, kb_id, f"{marker}{suffix}", content)
        file_ids.append(file_id)

    parsed_documents = await asyncio.gather(*(kb.parse_file(kb_id, file_id) for file_id in file_ids))
    for (suffix, marker), file_id, parsed in zip(formats, file_ids, parsed_documents, strict=True):
        assert parsed["status"] == "parsed"
        assert marker in parsed["parse_metadata"]["blocks"][0]["text"].upper()
        metadata = parsed["parse_metadata"]
        assert metadata["original_format"] == suffix.removeprefix(".")
        assert metadata["normalized_format"] == f"{suffix}x".removeprefix(".")
        assert metadata["conversion_required"] is True
        assert metadata["converter_name"] == "libreoffice"
        assert metadata["converter_version"].lower().startswith("libreoffice")
        assert metadata["conversion_duration_ms"] >= 0
        assert "conversion_warnings" in metadata

        indexed = await kb.index_file(kb_id, file_id)
        assert indexed["status"] == "indexed"
        record = await file_repository.get_by_file_id(file_id)
        assert record is not None
        assert record.chunk_count > 0
        assert record.token_count > 0
        chunks = await chunk_repository.list_by_file_id(file_id)
        assert chunks
        assert all(chunk.source_metadata["original_format"] == suffix.removeprefix(".") for chunk in chunks)
        assert all(chunk.source_metadata["normalized_format"] == f"{suffix}x".removeprefix(".") for chunk in chunks)
        results = await kb.aquery(marker, kb_id, search_mode="vector", final_top_k=20, similarity_threshold=-1.0)
        assert file_id in _result_file_ids(results)


async def test_gif_webp_real_first_frame_ocr_chunk_index_and_retrieve(
    test_client,
    admin_headers,
    knowledge_database,
    monkeypatch,
) -> None:
    kb_id = knowledge_database["kb_id"]
    documents = [
        (".gif", "GIFDELTA", _image_bytes("GIFDELTA OCR 48217", image_format="GIF", animated=True), 2),
        (".webp", "WEBPECHO", _image_bytes("WEBPECHO OCR 59328", image_format="WEBP"), 1),
    ]
    kb = await knowledge_base.aget_kb(kb_id)
    await _configure_deterministic_embedding(
        kb, kb_id, [marker for _suffix, marker, _content, _frames in documents], monkeypatch
    )
    file_repository = KnowledgeFileRepository()
    chunk_repository = KnowledgeChunkRepository()
    minio_client = get_minio_client()

    invalid = await test_client.post(
        "/api/knowledge/files/upload",
        params={"kb_id": kb_id},
        files={"file": ("broken.gif", b"GIF89a-broken", "image/gif")},
        headers=admin_headers,
    )
    assert invalid.status_code == 400

    for suffix, marker, content, frame_count in documents:
        file_id = await _upload_and_create(test_client, admin_headers, kb_id, f"{marker}{suffix}", content)
        record = await file_repository.get_by_file_id(file_id)
        assert record is not None
        source_bucket, source_object = parse_minio_url(record.path)
        assert await minio_client.astat_file(source_bucket, source_object) == len(content)

        parsed = await kb.parse_file(kb_id, file_id)
        metadata = parsed["parse_metadata"]
        assert parsed["status"] == "parsed"
        assert marker in metadata["blocks"][0]["text"].upper()
        assert metadata["parser_name"] == "rapid_ocr"
        assert metadata["original_format"] == suffix.removeprefix(".")
        assert metadata["normalized_format"] == "png"
        assert metadata["frame_count"] == frame_count
        assert metadata["selected_frames"] == [0]
        assert metadata["animation_ignored"] is (frame_count > 1)

        indexed = await kb.index_file(kb_id, file_id)
        assert indexed["status"] == "indexed"
        chunks = await chunk_repository.list_by_file_id(file_id)
        assert chunks
        assert all(chunk.source_metadata["parser_name"] == "rapid_ocr" for chunk in chunks)
        assert all(chunk.source_metadata["original_format"] == suffix.removeprefix(".") for chunk in chunks)
        assert all(chunk.source_metadata["normalized_format"] == "png" for chunk in chunks)
        assert all(chunk.source_metadata["page_number"] == 1 for chunk in chunks)
        results = await kb.aquery(marker, kb_id, search_mode="vector", final_top_k=20, similarity_threshold=-1.0)
        assert file_id in _result_file_ids(results)
