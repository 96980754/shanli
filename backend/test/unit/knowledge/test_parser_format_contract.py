from __future__ import annotations

import io
from pathlib import Path

import fitz
import pytest
import yuxi.knowledge.parser.unified as parser_unified
from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from yuxi.knowledge.parser.unified import Parser, validate_document_bytes


def _build_pdf(path: Path, page_texts: list[str]) -> None:
    document = fitz.open()
    for text in page_texts:
        page = document.new_page()
        if text:
            page.insert_text((72, 72), text)
    document.save(path)
    document.close()


def _build_docx(path: Path) -> None:
    document = Document()
    document.add_heading("DOCX title", level=1)
    document.add_paragraph("Paragraph before table")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "alpha"
    table.cell(1, 1).text = "one"
    document.add_paragraph("Paragraph after table")
    document.save(path)


def _build_xlsx(path: Path) -> None:
    workbook = Workbook()
    first = workbook.active
    first.title = "First"
    first.append(["Name", "Value"])
    first.append(["alpha", 1])
    second = workbook.create_sheet("Second")
    second.append(["Code", "Label"])
    second.append(["B", "beta"])
    workbook.create_sheet("Empty")
    workbook.save(path)


def _build_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Slide one"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    textbox.text = "First slide paragraph"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "alpha"
    table.cell(1, 1).text = "one"
    empty_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    assert empty_slide is not None
    presentation.save(path)


def test_txt_and_markdown_contract_preserves_content_and_bom(tmp_path: Path) -> None:
    txt_path = tmp_path / "plain.txt"
    txt_path.write_bytes("\ufeffPlain text content".encode())
    md_path = tmp_path / "notes.md"
    markdown = "# Heading\n\nParagraph with **formatting**.\n"
    md_path.write_text(markdown, encoding="utf-8")

    txt_result = Parser.parse_result(str(txt_path))
    md_result = Parser.parse_result(str(md_path))

    assert txt_result.markdown == "Plain text content"
    assert txt_result.parser_name == "native_text"
    assert txt_result.blocks[0].block_type == "paragraph"
    assert md_result.markdown == markdown
    assert md_result.document_title == "Heading"
    assert md_result.blocks[0].block_type == "heading"


@pytest.mark.parametrize("suffix", [".txt", ".md"])
def test_text_parser_rejects_blank_or_punctuation_only_documents(tmp_path: Path, suffix: str) -> None:
    path = tmp_path / f"empty{suffix}"
    path.write_text(" --- !!! \n", encoding="utf-8")

    with pytest.raises(ValueError, match="有效文本"):
        Parser.parse_result(str(path))


def test_text_parser_reports_utf8_decode_error(tmp_path: Path) -> None:
    path = tmp_path / "invalid.txt"
    path.write_bytes(b"\xff\xfe\x00\x00")

    with pytest.raises(ValueError, match="UTF-8"):
        Parser.parse_result(str(path))


def test_text_pdf_preserves_page_numbers_and_ignores_empty_pages(tmp_path: Path) -> None:
    path = tmp_path / "pages.pdf"
    _build_pdf(path, ["First page text", "", "Third page text"])

    result = Parser.parse_result(str(path))

    assert [block.page_number for block in result.blocks] == [1, 3]
    assert "First page text" in result.markdown
    assert "Third page text" in result.markdown
    assert result.parser_name == "native_pdf"


def test_scanned_pdf_fails_when_all_ocr_providers_fail(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    path = tmp_path / "scan.pdf"
    image = Image.new("RGB", (300, 120), "white")
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    document = fitz.open()
    page = document.new_page()
    page.insert_image(page.rect, stream=buffer.getvalue())
    document.save(path)
    document.close()

    async def fail_ocr(*_args, **_kwargs):
        from yuxi.knowledge.parser.ocr_routing import OCRRoutingError

        raise OCRRoutingError("OCR unavailable", attempts=[], warnings=["providers unavailable"])

    monkeypatch.setattr(parser_unified, "run_ocr_fallback", fail_ocr)

    with pytest.raises(ValueError, match="OCR"):
        Parser.parse_result(str(path))


def test_docx_fallback_preserves_paragraph_table_order(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    path = tmp_path / "ordered.docx"
    _build_docx(path)
    monkeypatch.setattr(
        parser_unified,
        "_convert_with_docling",
        lambda *args, **kwargs: (_ for _ in ()).throw(RuntimeError("force fallback")),
    )

    result = Parser.parse_result(str(path))

    block_types = [block.block_type for block in result.blocks]
    assert block_types == ["heading", "paragraph", "table", "paragraph"]
    assert result.markdown.index("Paragraph before table") < result.markdown.index("| Name | Value |")
    assert result.markdown.index("| Name | Value |") < result.markdown.index("Paragraph after table")
    assert result.parser_name == "python-docx"


def test_xlsx_preserves_sheet_order_and_simple_tables(tmp_path: Path) -> None:
    path = tmp_path / "book.xlsx"
    _build_xlsx(path)

    result = Parser.parse_result(str(path))

    assert [block.sheet_name for block in result.blocks] == ["First", "Second"]
    assert "## First" in result.markdown
    assert "| Name | Value |" in result.markdown
    assert result.markdown.index("## First") < result.markdown.index("## Second")
    assert "Empty" not in result.markdown


def test_pptx_preserves_slide_numbers_text_and_tables(tmp_path: Path) -> None:
    path = tmp_path / "slides.pptx"
    _build_pptx(path)

    result = Parser.parse_result(str(path))

    assert {block.slide_number for block in result.blocks} == {1}
    assert any(block.block_type == "table" for block in result.blocks)
    assert "Slide one" in result.markdown
    assert "First slide paragraph" in result.markdown
    assert "| Name | Value |" in result.markdown


@pytest.mark.parametrize(
    ("filename", "content", "message"),
    [
        ("fake.pdf", b"not a pdf", "PDF"),
        ("fake.docx", b"not a zip", "DOCX"),
        ("fake.xlsx", b"not a zip", "XLSX"),
        ("fake.pptx", b"not a zip", "PPTX"),
        ("legacy.doc", b"legacy", "OLE"),
    ],
)
def test_document_signature_validation_rejects_mismatch_or_unsupported(
    filename: str,
    content: bytes,
    message: str,
) -> None:
    with pytest.raises(ValueError, match=message):
        validate_document_bytes(filename, content)


def test_document_signature_validation_accepts_real_office_containers(tmp_path: Path) -> None:
    docx_path = tmp_path / "real.docx"
    xlsx_path = tmp_path / "real.xlsx"
    pptx_path = tmp_path / "real.pptx"
    _build_docx(docx_path)
    _build_xlsx(xlsx_path)
    _build_pptx(pptx_path)

    validate_document_bytes(docx_path.name, docx_path.read_bytes())
    validate_document_bytes(xlsx_path.name, xlsx_path.read_bytes())
    validate_document_bytes(pptx_path.name, pptx_path.read_bytes())


@pytest.mark.parametrize(("suffix", "format_name"), [(".png", "PNG"), (".jpg", "JPEG"), (".jpeg", "JPEG")])
def test_image_signature_validation_accepts_real_images(
    tmp_path: Path,
    suffix: str,
    format_name: str,
) -> None:
    path = tmp_path / f"real{suffix}"
    Image.new("RGB", (32, 24), "white").save(path, format=format_name)

    validate_document_bytes(path.name, path.read_bytes())


@pytest.mark.parametrize(
    ("filename", "content"),
    [
        ("fake.png", b"\x89PNG\r\n\x1a\nnot-an-image"),
        ("fake.jpg", b"\xff\xd8not-an-image\xff\xd9"),
        ("actually.png", b"\xff\xd8not-a-png\xff\xd9"),
    ],
)
def test_image_signature_validation_rejects_corruption_or_extension_mismatch(
    filename: str,
    content: bytes,
) -> None:
    with pytest.raises(ValueError):
        validate_document_bytes(filename, content)


def test_image_validation_rejects_excessive_pixel_count(tmp_path: Path) -> None:
    path = tmp_path / "large.png"
    Image.new("RGB", (11, 11), "white").save(path)

    with pytest.raises(ValueError, match="像素"):
        validate_document_bytes(path.name, path.read_bytes(), params={"ocr_max_image_pixels": 100})
